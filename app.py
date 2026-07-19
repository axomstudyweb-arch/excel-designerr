import os
import uuid
import logging
import io
import time
import threading
from pathlib import Path
from typing import Dict, Any, List, Optional
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
import pandas as pd
import openpyxl
from openpyxl.styles import PatternFill, Font
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

DATA_MUTATION_LOCK = threading.Lock()
SESSION_DELTAS: Dict[str, Dict[str, Any]] = {}
MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024 
WORKBOOK_CACHE: Dict[str, pd.DataFrame] = {}

class MutationPayload(BaseModel):
    data: Optional[Dict[str, Any]] = Field(default_factory=dict)
    headers_array: Optional[List[str]] = Field(default_factory=list)
    styles: Optional[Dict[str, Any]] = Field(default_factory=dict)
    title: Optional[str] = ""

class ExportPayload(BaseModel):
    theme: str
    title: Optional[str] = ""
    total_rows: int

def run_storage_cleaner():
    while True:
        try:
            now = time.time()
            for item in UPLOAD_DIR.iterdir():
                if item.is_file() and (now - item.stat().st_mtime) > 3600:
                    item.unlink()
                    wb_id = item.stem
                    with DATA_MUTATION_LOCK:
                        SESSION_DELTAS.pop(wb_id, None)
                        WORKBOOK_CACHE.pop(wb_id, None)
        except Exception as e:
            logger.error(f"Cleaner thread error: {e}")
        time.sleep(600)

threading.Thread(target=run_storage_cleaner, daemon=True).start()

class HighPerformanceProcessor:
    def parse_csv_with_encoding_fallback(self, file_path: Path) -> pd.DataFrame:
        for enc in ['utf-8', 'utf-8-sig', 'cp1252', 'latin1']:
            try: return pd.read_csv(file_path, header=None, encoding=enc)
            except Exception: continue
        raise ValueError("Unsupported matrix character structure mapping.")

    def find_header_offset(self, df: pd.DataFrame) -> int:
        keywords = {
            'si', 'rc', 'id', 'date', 'name', 'serial', 'no', 'code', 'sl no', 
            'beneficiary name', 'father name', 'bank name', 'ifsc', 'account number', 
            'aadhaar number', 'aadhar', 'mobile', 'state', 'district', 'gp'
        }
        try:
            for idx, row in df.head(25).iterrows():
                row_vals = [str(x).strip().lower() for x in row.values if pd.notna(x)]
                if any(k in row_vals for k in keywords): return idx
        except Exception: pass
        return 0

processor = HighPerformanceProcessor()

@app.get("/")
async def serve_index():
    if os.path.exists("templates/index.html"): return FileResponse("templates/index.html")
    if os.path.exists("index.html"): return FileResponse("index.html")
    raise HTTPException(status_code=404, detail="index.html template layout missing.")

@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    suffix = Path(file.filename).suffix.lower()
    if suffix not in {".xlsx", ".xls", ".csv"}:
        raise HTTPException(status_code=400, detail="Invalid extension format.")
        
    wb_id = str(uuid.uuid4())
    temp_path = UPLOAD_DIR / f"{wb_id}{suffix}"
    
    total_bytes_written = 0
    try:
        with open(temp_path, "wb") as buffer:
            while chunk := await file.read(512 * 1024):
                total_bytes_written += len(chunk)
                if total_bytes_written > MAX_FILE_SIZE_BYTES:
                    raise HTTPException(status_code=413, detail="Payload content length limits exceeded.")
                buffer.write(chunk)
    except Exception as e:
        if temp_path.exists(): temp_path.unlink()
        raise e

    try:
        if suffix == '.csv': df = processor.parse_csv_with_encoding_fallback(temp_path)
        else: df = pd.read_excel(temp_path, header=None, engine='openpyxl')
        WORKBOOK_CACHE[wb_id] = df
    except Exception as err:
        if temp_path.exists(): temp_path.unlink()
        raise HTTPException(status_code=500, detail=f"Parsing failure: {str(err)}")

    offset = processor.find_header_offset(df)
    total_rows = len(df)
    
    try:
        raw_h = df.iloc[offset].values.tolist()
        headers = []
        seen = {}
        for i, x in enumerate(raw_h):
            base_name = str(x).strip() if pd.notna(x) else f"Column {i+1}"
            if base_name in seen:
                seen[base_name] += 1
                headers.append(f"{base_name}_{seen[base_name]}")
            else:
                seen[base_name] = 0
                headers.append(base_name)
    except Exception:
        headers = [f"Column {i+1}" for i in range(df.shape[1])]

    with DATA_MUTATION_LOCK:
        # Save the structural header offset configuration to correctly parse rows later
        SESSION_DELTAS[wb_id] = {"data": {}, "headers": headers, "styles": {}, "title": "", "offset": offset}

    return {
        "workbook_id": wb_id,
        "meta": {
            "rows": max(1, total_rows - offset - 1),
            "cols": len(headers),
            "offset": offset,
            "headers": headers
        }
    }

@app.get("/api/workbook/{workbook_id}/rows")
async def get_rows(workbook_id: str, start: int, count: int, offset: int):
    if workbook_id not in WORKBOOK_CACHE:
        raise HTTPException(status_code=404, detail="Active workbook context footprint missing.")
    df = WORKBOOK_CACHE[workbook_id]
    adjusted_start = offset + 1 + start
    try:
        slice_df = df.iloc[adjusted_start : adjusted_start + count]
        return {"rows": slice_df.fillna("").values.tolist(), "start": start}
    except Exception:
        return {"rows": [], "start": start}

@app.post("/api/mutate/{workbook_id}")
async def mutate_session(workbook_id: str, payload: MutationPayload):
    with DATA_MUTATION_LOCK:
        if workbook_id not in SESSION_DELTAS: raise HTTPException(status_code=404, detail="Session expired.")
        if payload.data: SESSION_DELTAS[workbook_id]["data"].update(payload.data)
        if payload.headers_array: SESSION_DELTAS[workbook_id]["headers"] = payload.headers_array
        if payload.styles: SESSION_DELTAS[workbook_id]["styles"].update(payload.styles)
        if payload.title is not None: SESSION_DELTAS[workbook_id]["title"] = payload.title
    return {"status": "synced"}

@app.post("/api/export/{workbook_id}")
async def export_workbook(workbook_id: str, payload: ExportPayload):
    with DATA_MUTATION_LOCK:
        if workbook_id not in SESSION_DELTAS: raise HTTPException(status_code=404, detail="Session expired.")
        session = SESSION_DELTAS[workbook_id]
        df = WORKBOOK_CACHE.get(workbook_id)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Designer Export"

    excel_theme_fills = {
        "progress": {"header": "F8FAFC", "even": "FFFFFF", "odd": "FDFBF7"},
        "teal-dashboard": {"header": "005F56", "even": "FFFFFF", "odd": "FAFAFA"},
        "financial": {"header": "104F60", "even": "E8EEF3", "odd": "FFFFFF"},
        "emerald": {"header": "107C41", "even": "FFFFFF", "odd": "F8FAFC"}
    }
    t_set = excel_theme_fills.get(payload.theme, excel_theme_fills["emerald"])

    headers = session["headers"]
    for c_idx, h_val in enumerate(headers):
        cell = ws.cell(row=1, column=c_idx + 1, value=str(h_val))
        cell.fill = PatternFill(start_color=t_set["header"], end_color=t_set["header"], fill_type="solid")
        if payload.theme in ["teal-dashboard", "financial", "emerald"]:
            cell.font = Font(color="FFFFFF", bold=True)

    offset = session.get("offset", 0)

    for r in range(1, payload.total_rows + 1):
        row_color = t_set["even"] if (r % 2 == 0) else t_set["odd"]
        for c in range(len(headers)):
            # 1st Fallback Strategy: Check if the cell has frontend modifications
            cell_val = session["data"].get(f"{r}-{c}", session["data"].get(f"{r}-${c}", None))
            
            # 2nd Fallback Strategy: Read straight from original DataFrame if not touched on frontend
            if cell_val is None and df is not None:
                df_row_idx = offset + r
                if df_row_idx < len(df) and c < df.shape[1]:
                    val_raw = df.iloc[df_row_idx, c]
                    cell_val = str(val_raw) if pd.notna(val_raw) else ""
                else:
                    cell_val = ""
            elif cell_val is None:
                cell_val = ""

            # Explicit Privacy Control Layer: Mask sensitive government unique indicators
            h_name = str(headers[c]).lower() if c < len(headers) else ""
            if "aadhar" in h_name or "aadhaar" in h_name:
                cell_val = "[Aadhaar Redacted]"

            cell = ws.cell(row=r + 1, column=c + 1, value=str(cell_val))
            
            cell_style = session["styles"].get(f"{r}-{c}", session["styles"].get(f"{r}-${c}", {}))
            custom_tint = cell_style.get("tint")
            
            if isinstance(custom_tint, str) and custom_tint.startswith("#"):
                hex_color = custom_tint.lstrip("#").upper()
                cell.fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")
            else:
                cell.fill = PatternFill(start_color=row_color, end_color=row_color, fill_type="solid")
                
            if cell_style.get("bold"):
                cell.font = Font(bold=True)

    # Append Signatory Block at structural bottom context
    start_sig_row = payload.total_rows + 3
    signatory_data = [
        ("Software", "Excel Engine Designer"),
        ("Version", "2.0"),
        ("Author", "Ankur Dowarah"),
        ("URL", "ank-ldb.blogspot.com"),
        ("Copyright", "Copyright © 2026 Ankur Dowarah")
    ]
    
    for idx, (label, val) in enumerate(signatory_data):
        curr_row = start_sig_row + idx
        c_lbl = ws.cell(row=curr_row, column=1, value=f"{label}:")
        c_lbl.font = Font(bold=True, color="475569")
        c_val = ws.cell(row=curr_row, column=2, value=val)
        c_val.font = Font(color="0F172A")

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return StreamingResponse(out, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.post("/api/export_pptx/{workbook_id}")
async def export_pptx(workbook_id: str, payload: ExportPayload):
    with DATA_MUTATION_LOCK:
        if workbook_id not in SESSION_DELTAS: raise HTTPException(status_code=404, detail="Session expired.")
        session = SESSION_DELTAS[workbook_id]
        df = WORKBOOK_CACHE.get(workbook_id)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 

    if payload.title:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = payload.title
        p.font.size = Pt(22)
        p.font.bold = True

    headers = session["headers"]
    rows_count = min(payload.total_rows, 10) 
    cols_count = len(headers)

    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9.0)
    height = Inches(0.35 * (rows_count + 1))

    table_shape = slide.shapes.add_table(rows_count + 1, cols_count, left, top, width, height)
    table = table_shape.table

    for c_idx, h_val in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = str(h_val)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(16, 124, 65)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    offset = session.get("offset", 0)

    for r in range(1, rows_count + 1):
        for c in range(cols_count):
            # Fallback evaluation sequence logic for PowerPoint compilation matrix
            cell_val = session["data"].get(f"{r}-{c}", session["data"].get(f"{r}-${c}", None))
            if cell_val is None and df is not None:
                df_row_idx = offset + r
                if df_row_idx < len(df) and c < df.shape[1]:
                    val_raw = df.iloc[df_row_idx, c]
                    cell_val = str(val_raw) if pd.notna(val_raw) else ""
                else:
                    cell_val = ""
            elif cell_val is None:
                cell_val = ""

            h_name = str(headers[c]).lower() if c < len(headers) else ""
            if "aadhar" in h_name or "aadhaar" in h_name:
                cell_val = "[Aadhaar Redacted]"

            cell = table.cell(r, c)
            cell.text = cell_val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)

    # Append Signature Block Context below PPTX table shape
    sig_box_top = top + height + Inches(0.3)
    sigBox = slide.shapes.add_textbox(Inches(0.5), sig_box_top, Inches(5), Inches(1.2))
    tf_sig = sigBox.text_frame
    tf_sig.word_wrap = True
    
    signatory_text = (
        "Software: Excel Engine Designer\n"
        "Version: 2.0\n"
        "Author: Ankur Dowarah\n"
        "URL: ank-ldb.blogspot.com\n"
        "Copyright © 2026 Ankur Dowarah"
    )
    
    p_sig = tf_sig.paragraphs[0]
    p_sig.text = signatory_text
    p_sig.font.size = Pt(10)
    p_sig.font.color.rgb = RGBColor(71, 85, 105)
    p_sig.font.italic = True

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return StreamingResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")